import shutil
import subprocess
from io import BytesIO
from pathlib import Path
from textwrap import wrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
PPTX_PATH = OUT_DIR / "latent_gepa_discussione.pptx"
PREVIEW_DIR = OUT_DIR / "preview"
QA_PATH = OUT_DIR / "deck_qa_report.md"
FIGURE_ASSET_DIR = OUT_DIR / "figure_assets"
THESIS_FIGURE_DIR = OUT_DIR.parent / "thesis" / "latex" / "figures"

WIDE_W = 13.333333
WIDE_H = 7.5

COLORS = {
    "black": RGBColor(18, 24, 32),
    "gray": RGBColor(88, 96, 105),
    "light_gray": RGBColor(242, 244, 247),
    "mid_gray": RGBColor(205, 211, 219),
    "blue": RGBColor(24, 79, 138),
    "blue_light": RGBColor(226, 238, 250),
    "green": RGBColor(34, 126, 85),
    "green_light": RGBColor(227, 244, 235),
    "red": RGBColor(174, 60, 60),
    "red_light": RGBColor(248, 228, 228),
    "orange": RGBColor(201, 123, 38),
    "orange_light": RGBColor(252, 238, 218),
    "purple": RGBColor(101, 71, 160),
    "purple_light": RGBColor(238, 232, 250),
    "white": RGBColor(255, 255, 255),
}

FONT_SCALE = (8.5, 11.0, 13.5, 16.0, 20.0, 28.0, 34.0)

TITLE = (
    "Latent-GEPA: Accelerating Prompt Optimization via Embedding Inversion "
    "and Latent Semantic Guidance in Large Language Models"
)

SLIDE_TITLES = [
    "Title",
    "Why Latent Text Needs Semantic Checks",
    "From Readouts to Prompt Optimization",
    "Thesis Contributions",
    "Semantic-Fidelity Corpus",
    "Latent-GEPA Feedback Loop",
    "What Feedback Was Tested",
    "Useful Does Not Mean Interpretable",
    "Perplexity Helps Most Clearly",
    "NLA Is Diagnostic, Not Yet Conclusive",
    "Conclusions and Future Work",
]


def rgb(name):
    return COLORS[name]


def normalized_font_size(size):
    if size <= 9:
        return FONT_SCALE[0]
    if size <= 12.5:
        return FONT_SCALE[1]
    if size <= 16.5:
        return FONT_SCALE[2]
    if size <= 19:
        return FONT_SCALE[3]
    if size <= 23:
        return FONT_SCALE[4]
    if size <= 30:
        return FONT_SCALE[5]
    return FONT_SCALE[6]


def set_font(run, size=18, bold=False, color="black", name="Aptos"):
    run.font.name = name
    run.font.size = Pt(normalized_font_size(size))
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def textbox(slide, x, y, w, h, text="", size=18, color="black", bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if text else [""]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = line
        set_font(r, size=size, bold=bold, color=color)
    return shape


def box(
    slide,
    x,
    y,
    w,
    h,
    text,
    fill="light_gray",
    line="mid_gray",
    size=16,
    bold=False,
    color="black",
    align=PP_ALIGN.CENTER,
    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = line
        set_font(r, size=size, bold=bold or idx == 0, color=color)
    return shape


def title(slide, text, accent=None):
    textbox(slide, 0.55, 0.28, 12.25, 0.55, text, size=34, bold=True, color="black")
    if accent:
        textbox(slide, 0.57, 0.84, 4.0, 0.18, accent, size=9, color="gray")


def footer(slide, number):
    textbox(slide, 0.55, 7.12, 5.0, 0.22, "Simone Mazzacano | Latent-GEPA", size=8.5, color="gray")
    textbox(slide, 5.45, 7.12, 2.4, 0.22, "Bachelor Thesis Defense", size=8.5, color="gray", align=PP_ALIGN.CENTER)
    textbox(slide, 12.35, 7.12, 0.45, 0.22, str(number), size=8.5, color="gray", align=PP_ALIGN.RIGHT)


def add_link(slide, x, y, label, url):
    shape = textbox(slide, x, y, 4.35, 0.20, "", size=8.5, color="blue")
    p = shape.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = label
    set_font(r, size=8.5, color="blue")
    r.hyperlink.address = url
    return shape


def chip(slide, x, y, text, fill, line=None, w=None):
    width = w if w is not None else max(0.95, 0.12 * len(text) + 0.28)
    return box(slide, x, y, width, 0.34, text, fill=fill, line=line or fill, size=11, bold=True)


def arrow(slide, x, y, w, h=0.28, color="blue_light", line="blue"):
    return box(
        slide,
        x,
        y,
        w,
        h,
        "",
        fill=color,
        line=line,
        shape_type=MSO_SHAPE.RIGHT_ARROW,
    )


def ensure_pdf_figure(name, source):
    FIGURE_ASSET_DIR.mkdir(exist_ok=True)
    pdf_path = THESIS_FIGURE_DIR / source
    png_path = FIGURE_ASSET_DIR / f"{name}.png"
    if not pdf_path.exists():
        raise FileNotFoundError(f"Missing thesis figure source: {pdf_path}")
    if png_path.exists() and png_path.stat().st_mtime >= pdf_path.stat().st_mtime:
        return png_path
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        raise RuntimeError("pdftoppm is required to convert thesis PDF figures to PPTX images")
    subprocess.run(
        [pdftoppm, "-singlefile", "-png", "-r", "220", str(pdf_path), str(png_path.with_suffix(""))],
        check=True,
    )
    return png_path


def draw_centered_lines(draw, center_x, top_y, lines, font, fill, spacing=4):
    y = top_y
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=font)
        draw.text((center_x - (bbox[2] - bbox[0]) / 2, y), line, font=font, fill=fill)
        y += bbox[3] - bbox[1] + spacing


def create_soft_prompt_deltas_figure():
    FIGURE_ASSET_DIR.mkdir(exist_ok=True)
    png_path = FIGURE_ASSET_DIR / "soft_prompt_deltas_labeled.png"
    labels = [
        ("Main", "setting"),
        ("Fewer", "virtual tokens"),
        ("More", "virtual tokens"),
        ("Shorter", "context"),
        ("Random", "seed 43"),
        ("Random", "seed 44"),
    ]
    validation = [0.0313, 0.0047, 0.0431, 0.0277, -0.0733, -0.0632]
    final_test = [0.0557, 0.0218, -0.0384, 0.0233, 0.0863, -0.1008]

    width, height = 1500, 760
    margin_left, margin_right = 120, 45
    margin_top, margin_bottom = 95, 185
    plot_w = width - margin_left - margin_right
    plot_h = height - margin_top - margin_bottom
    y_min, y_max = -0.12, 0.10

    img = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(img)
    title_font = get_font(34, bold=True)
    axis_font = get_font(24)
    tick_font = get_font(22)
    label_font = get_font(20, bold=True)
    small_font = get_font(20)

    def sx(group_idx, offset=0):
        group_w = plot_w / len(labels)
        return margin_left + group_idx * group_w + group_w / 2 + offset

    def sy(value):
        return margin_top + (y_max - value) / (y_max - y_min) * plot_h

    draw.text((margin_left, 20), "Soft-prompt task deltas by variant", font=title_font, fill=(18, 24, 32))

    legend_y = 35
    draw.rectangle([880, legend_y, 920, legend_y + 20], fill=(52, 102, 153))
    draw.text((935, legend_y - 6), "Validation", font=axis_font, fill=(18, 24, 32))
    draw.rectangle([1095, legend_y, 1135, legend_y + 20], fill=(203, 132, 39))
    draw.text((1150, legend_y - 6), "Final test", font=axis_font, fill=(18, 24, 32))

    for tick in [-0.10, -0.05, 0.00, 0.05, 0.10]:
        y = sy(tick)
        color = (18, 24, 32) if tick == 0 else (216, 222, 230)
        draw.line([margin_left, y, width - margin_right, y], fill=color, width=3 if tick == 0 else 1)
        tick_label = f"{tick:.2f}"
        bbox = draw.textbbox((0, 0), tick_label, font=tick_font)
        draw.text((margin_left - bbox[2] - 16, y - 13), tick_label, font=tick_font, fill=(18, 24, 32))

    draw.line([margin_left, margin_top, margin_left, margin_top + plot_h], fill=(18, 24, 32), width=2)
    draw.text((margin_left, margin_top - 38), "Pearson delta", font=axis_font, fill=(18, 24, 32))

    group_w = plot_w / len(labels)
    bar_w = 34
    zero_y = sy(0)
    for idx, (val, test) in enumerate(zip(validation, final_test)):
        cx = sx(idx)
        for offset, value, color in [(-bar_w / 1.7, val, (52, 102, 153)), (bar_w / 1.7, test, (203, 132, 39))]:
            x0 = cx + offset - bar_w / 2
            x1 = cx + offset + bar_w / 2
            y = sy(value)
            draw.rectangle([x0, min(y, zero_y), x1, max(y, zero_y)], fill=color)
        draw_centered_lines(draw, cx, margin_top + plot_h + 24, labels[idx], label_font, (18, 24, 32), spacing=3)

    draw.text(
        (margin_left + plot_w / 2 - 350, height - 42),
        "Main setting: 16 virtual tokens, 2048-token context, random seed 42.",
        font=small_font,
        fill=(88, 96, 105),
    )
    img.save(png_path)
    return png_path


def add_figure(slide, x, y, w, h, image_path, line="mid_gray"):
    box(slide, x, y, w, h, "", fill="white", line=line, shape_type=MSO_SHAPE.RECTANGLE)
    return slide.shapes.add_picture(
        str(image_path),
        Inches(x + 0.06),
        Inches(y + 0.06),
        width=Inches(w - 0.12),
        height=Inches(h - 0.12),
    )


def bullet_panel(slide, x, y, w, h, heading, bullets, fill="white", line="mid_gray", heading_color="blue"):
    shape = box(slide, x, y, w, h, heading, fill=fill, line=line, size=15, bold=True, color=heading_color, align=PP_ALIGN.LEFT)
    tf = shape.text_frame
    for bullet in bullets:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = bullet
        set_font(r, size=12.5, color="black")
    return shape


def metric_card(slide, x, y, label, before, after, good=True):
    fill = "green_light" if good else "red_light"
    line = "green" if good else "red"
    shape = box(slide, x, y, 2.75, 1.12, label, fill=fill, line=line, size=15, bold=True, align=PP_ALIGN.LEFT)
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = f"{before} -> {after}"
    set_font(r, size=20, bold=True, color=line)
    return shape


def bar_group(slide, x, y, label, values, colors, max_value=1.0):
    textbox(slide, x, y - 0.22, 2.3, 0.2, label, size=10.5, color="gray", bold=True, align=PP_ALIGN.CENTER)
    base_y = y + 1.15
    for idx, (v, c) in enumerate(zip(values, colors)):
        bh = max(0.08, 1.05 * v / max_value)
        bx = x + 0.18 + idx * 0.43
        box(slide, bx, base_y - bh, 0.28, bh, "", fill=c, line=c, shape_type=MSO_SHAPE.RECTANGLE)
        textbox(slide, bx - 0.07, base_y + 0.03, 0.42, 0.16, str(idx + 1), size=7.5, color="gray", align=PP_ALIGN.CENTER)


def build_deck():
    soft_prompt_deltas = create_soft_prompt_deltas_figure()
    gepa_long_run_metrics = ensure_pdf_figure("gepa_long_run_metrics", Path("chapter5/gepa_long_run_metrics.pdf"))

    prs = Presentation()
    prs.slide_width = Inches(WIDE_W)
    prs.slide_height = Inches(WIDE_H)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    textbox(slide, 0.72, 0.72, 11.9, 1.55, TITLE, size=28, bold=True, color="black", align=PP_ALIGN.CENTER)
    box(slide, 3.65, 2.55, 6.0, 0.48, "Bachelor Thesis in Data Intensive Applications", fill="blue_light", line="blue", size=16, bold=True)
    textbox(slide, 3.2, 3.38, 6.9, 0.36, "Candidate: Simone Mazzacano", size=19, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 3.2, 3.88, 6.9, 0.30, "Supervisor: Prof. Gianluca Moro", size=15, align=PP_ALIGN.CENTER)
    textbox(slide, 2.5, 4.25, 8.3, 0.34, "Co-supervisors: Dr. Lorenzo Molfetta, Dr. Stefano Fantazzini, Dr. Giacomo Frisoni", size=14, align=PP_ALIGN.CENTER)
    textbox(slide, 3.2, 5.05, 6.9, 0.34, "Computer Science and Engineering | Academic Year 2025-2026 | First Session", size=12.5, color="gray", align=PP_ALIGN.CENTER)
    for idx, (x, y, c) in enumerate([(1.4, 5.8, "blue_light"), (2.15, 5.45, "green_light"), (10.4, 5.72, "orange_light"), (11.15, 5.33, "purple_light")]):
        box(slide, x, y, 0.58, 0.58, "", fill=c, line="mid_gray", shape_type=MSO_SHAPE.OVAL)
        if idx < 2:
            arrow(slide, x + 0.56, y + 0.18, 0.82, 0.17, color="light_gray", line="mid_gray")

    # Slide 2
    slide = prs.slides.add_slide(blank)
    title(slide, SLIDE_TITLES[1])
    panels = [
        ("Problem", "Large Language Model (LLM)\nlatent states can be decoded\ninto readable text", "blue_light", "blue"),
        ("Why it matters", "Fluent reconstructions\ncan lose negation or\ncommonsense meaning", "orange_light", "orange"),
        ("New in this thesis", "Semantic-Fidelity Corpus\nLatent-GEPA tests\nsoft-prompt diagnostic", "green_light", "green"),
        ("Roadmap", "Semantic stress tests\n-> feedback loop\n-> matched evidence", "purple_light", "purple"),
    ]
    for x, (head, body, fill, line) in zip([0.68, 3.82, 6.96, 10.10], panels):
        box(slide, x, 1.18, 2.55, 1.28, f"{head}\n{body}", fill=fill, line=line, size=12.2, bold=True)
    textbox(slide, 0.72, 2.95, 11.75, 0.28, "Concrete semantic-drift example", size=15, bold=True, color="blue", align=PP_ALIGN.CENTER)
    box(slide, 1.10, 3.48, 4.85, 0.74, "The ice does not melt under the sun", fill="green_light", line="green", size=20, bold=True)
    arrow(slide, 6.20, 3.66, 0.82, 0.22, color="light_gray", line="mid_gray")
    box(slide, 7.34, 3.48, 4.85, 0.74, "The ice melts under the sun", fill="red_light", line="red", size=20, bold=True)
    chip(slide, 3.00, 4.48, "Lexically close", "light_gray", "mid_gray", w=1.85)
    chip(slide, 8.65, 4.48, "Negation removed", "red_light", "red", w=2.05)
    box(slide, 1.28, 5.42, 10.72, 0.68, "Surface overlap is not semantic fidelity\nTest meaning first; then test whether latent evidence improves judge-prompt search", fill="white", line="blue", size=16.5, bold=True)
    footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    title(slide, SLIDE_TITLES[2])
    steps = [
        ("Task\nexample", "blue_light", "blue"),
        ("Base judge\nfixed weights", "blue_light", "blue"),
        ("Score +\nerror", "orange_light", "orange"),
        ("Feedback\nrecord", "orange_light", "orange"),
        ("Proposer\nrewrites prompt", "green_light", "green"),
        ("New judge\nprompt", "green_light", "green"),
    ]
    x = 0.65
    for idx, (text, fill, line) in enumerate(steps):
        box(slide, x, 1.55, 1.68, 1.0, text, fill=fill, line=line, size=14, bold=True)
        if idx < len(steps) - 1:
            arrow(slide, x + 1.72, 1.90, 0.50, 0.23, color="light_gray", line="mid_gray")
        x += 2.05
    bullet_panel(
        slide,
        0.78,
        3.05,
        5.72,
        2.08,
        "LLM-as-a-judge",
        [
            "The model scores generated text against a rubric",
            "Quality is measured by agreement with human judgments",
            "The evaluated judge remains fixed during prompt search",
        ],
        fill="white",
        line="blue",
    )
    bullet_panel(
        slide,
        6.82,
        3.05,
        5.72,
        2.08,
        "Latent-GEPA idea",
        [
            "GEPA: reflective prompt search",
            "Search space: natural-language prompts",
            "Latent signals are attached to reflection feedback",
            "The proposer receives richer evidence for prompt revision",
        ],
        fill="white",
        line="green",
        heading_color="green",
    )
    textbox(slide, 0.95, 5.55, 11.25, 0.38, "Key terms: Base judge | Proposer | Feedback | Held-out evaluation", size=16, bold=True, color="blue", align=PP_ALIGN.CENTER)
    add_link(slide, 0.78, 6.30, "GEPA: Agrawal et al., 2025", "https://doi.org/10.48550/arXiv.2507.19457")
    add_link(slide, 4.30, 6.30, "G-EVAL: Liu et al., 2023", "https://doi.org/10.18653/v1/2023.emnlp-main.153")
    footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    title(slide, SLIDE_TITLES[3])
    contribs = [
        ("RQ1\nSemantic fidelity", "New controlled corpus for fragile meaning", "Corpus validated", "blue_light", "blue"),
        ("RQ2\nLatent feedback", "GEPA feedback enriched with PPL and activation text", "PPL helps most clearly", "green_light", "green"),
        ("RQ3\nSoft prompts", "Continuous-prompt readout tests semantic transparency", "Useful but opaque", "red_light", "red"),
    ]
    xs = [0.75, 4.66, 8.57]
    for x, (head, body, result, fill, line) in zip(xs, contribs):
        box(slide, x, 1.45, 3.35, 0.86, head, fill=fill, line=line, size=18, bold=True)
        box(slide, x, 2.58, 3.35, 1.30, body, fill="white", line=line, size=15)
        box(slide, x, 4.15, 3.35, 0.72, result, fill=fill, line=line, size=15, bold=True)
    textbox(slide, 0.88, 5.55, 11.35, 0.48, "The contributions connect the corpus, the feedback loop, and the final claim boundaries", size=18, bold=True, color="blue", align=PP_ALIGN.CENTER)
    footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    title(slide, SLIDE_TITLES[4])
    textbox(slide, 0.72, 1.02, 11.8, 0.34, "A new 2080-row surface: 1060 examples or contrastive pairs built around meaning-changing contrasts", size=17, bold=True, color="blue", align=PP_ALIGN.CENTER)
    blocks = [
        ("A", "40 rows", "Standard controls", "Ordinary short sentences", "blue_light", "blue"),
        ("B", "720 rows", "Negation and polarity", "Positive/negative pairs", "red_light", "red"),
        ("C", "1320 rows", "Commonsense stress", "Plausible vs counterfactual pairs", "green_light", "green"),
    ]
    xs = [0.75, 4.70, 8.65]
    for x, (letter, rows, name, desc, fill, line) in zip(xs, blocks):
        box(slide, x, 1.65, 3.15, 0.62, f"Block {letter}", fill=fill, line=line, size=18, bold=True)
        box(slide, x, 2.42, 3.15, 1.22, f"{rows}\n{name}", fill="white", line=line, size=17, bold=True)
        box(slide, x, 3.88, 3.15, 0.76, desc, fill=fill, line=line, size=13.5)
    examples = [
        ("A example", "The meeting starts\nat nine in the morning", "blue_light", "blue"),
        ("B contrast", "The church is empty of song\n-> filled with song", "red_light", "red"),
        ("C contrast", "when it rains humidity forms\n-> when it is hot humidity forms", "green_light", "green"),
    ]
    for x, (head, body, fill, line) in zip(xs, examples):
        box(slide, x, 5.00, 3.15, 0.92, f"{head}\n{body}", fill=fill, line=line, size=12.4, bold=True)
    textbox(slide, 1.05, 6.27, 11.10, 0.34, "Stable metadata joins inversion, activation verbalization, and scoring outputs", size=16, bold=True, color="gray", align=PP_ALIGN.CENTER)
    footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    title(slide, SLIDE_TITLES[5])
    nodes = [
        ("Seed\nprompt", "green_light", "green"),
        ("Task\nexamples", "blue_light", "blue"),
        ("Base judge\nevaluation", "blue_light", "blue"),
        ("Reflection\nfeedback", "orange_light", "orange"),
        ("Proposer\ncandidate prompts", "green_light", "green"),
        ("Validation\nfrontier", "green_light", "green"),
    ]
    x = 0.55
    for idx, (text, fill, line) in enumerate(nodes):
        box(slide, x, 1.70, 1.72, 0.92, text, fill=fill, line=line, size=13.5, bold=True)
        if idx < len(nodes) - 1:
            arrow(slide, x + 1.76, 2.02, 0.36, 0.18, color="light_gray", line="mid_gray")
        x += 2.05
    box(slide, 3.58, 3.30, 5.92, 1.30, "Latent feedback attached to reflection data\nPerplexity (PPL): response surprisal\nNatural Language Activations (NLA): activation text", fill="orange_light", line="orange", size=16.5, bold=True)
    chip(slide, 4.32, 4.95, "Metric", "blue_light", "blue", w=1.12)
    chip(slide, 5.70, 4.95, "PPL", "green_light", "green", w=0.82)
    chip(slide, 6.80, 4.95, "NLA", "red_light", "red", w=0.82)
    chip(slide, 7.90, 4.95, "Optional aux compression", "purple_light", "purple", w=2.42)
    textbox(slide, 0.78, 5.95, 11.75, 0.52, "Final held-out evaluation is run only after validation-time prompt selection", size=18, bold=True, color="blue", align=PP_ALIGN.CENTER)
    add_link(slide, 0.78, 6.48, "NLA: Fraser-Taliente et al., 2026", "https://transformer-circuits.pub/2026/nla/")
    footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    title(slide, SLIDE_TITLES[6])
    labels = [
        ("Metric feedback", "Target, prediction, and error summary", "blue_light", "blue"),
        ("Metric + PPL", "Adds response-only perplexity from the base judge", "green_light", "green"),
        ("Metric + PPL + raw NLA", "Adds verbalizations from selected token positions", "red_light", "red"),
        ("Metric + PPL + schema-fixed NLA", "Cleaner metadata and token status; verbalizer not retrained", "purple_light", "purple"),
    ]
    y = 1.22
    for name, desc, fill, line in labels:
        box(slide, 0.78, y, 3.25, 0.62, name, fill=fill, line=line, size=14.5, bold=True)
        box(slide, 4.28, y, 7.75, 0.62, desc, fill="white", line=line, size=14.2, align=PP_ALIGN.LEFT)
        y += 1.02
    bullet_panel(
        slide,
        0.95,
        5.55,
        5.15,
        0.98,
        "Main evaluation task",
        ["Topical-Chat / USR engagingness", "Human agreement metrics on held-out rows"],
        fill="white",
        line="blue",
    )
    bullet_panel(
        slide,
        7.05,
        5.55,
        5.15,
        0.98,
        "Model roles",
        ["Qwen2.5-7B judge/PPL model; layer-20 verbalizer", "Qwen35B proposer; final scores from 7B judge"],
        fill="white",
        line="green",
        heading_color="green",
    )
    footer(slide, 7)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    title(slide, SLIDE_TITLES[7])
    box(slide, 0.78, 1.25, 3.45, 1.12, "Discrete tokens\nSIPIT exact recovery\n100/100 GPT-2; 40/40 logical", fill="green_light", line="green", size=14.2, bold=True, align=PP_ALIGN.LEFT)
    box(slide, 4.48, 1.25, 3.45, 1.12, "Embedding inversion\nFull-mask probes weak\nTiny overfit positive control", fill="orange_light", line="orange", size=14.2, bold=True, align=PP_ALIGN.LEFT)
    box(slide, 8.18, 1.25, 3.45, 1.12, "Soft prompts\nTask-active runs exist\nExact readout still false", fill="red_light", line="red", size=14.2, bold=True, align=PP_ALIGN.LEFT)
    add_figure(slide, 0.88, 2.70, 7.55, 3.48, soft_prompt_deltas, line="red")
    box(slide, 8.78, 2.82, 3.55, 0.80, "Relabeled source\nsoft_prompt_deltas.pdf", fill="light_gray", line="mid_gray", size=14, bold=True)
    box(slide, 8.78, 3.92, 3.55, 1.04, "The positive soft-prompt deltas do not make the learned vectors semantically transparent", fill="red_light", line="red", size=14.8, bold=True)
    box(slide, 8.78, 5.22, 3.55, 0.78, "Boundary: usefulness and interpretability are separate claims", fill="blue_light", line="blue", size=14.6, bold=True)
    add_link(slide, 0.78, 6.48, "SIPIT: Nikolaou et al., 2025", "https://doi.org/10.48550/arXiv.2510.15511")
    footer(slide, 8)

    # Slide 9
    slide = prs.slides.add_slide(blank)
    title(slide, SLIDE_TITLES[8])
    metric_card(slide, 0.78, 1.30, "Pearson", "0.5518", "0.6328")
    metric_card(slide, 3.72, 1.30, "Spearman", "0.5478", "0.6199")
    metric_card(slide, 6.66, 1.30, "Agreement", "0.7611", "0.7889")
    metric_card(slide, 9.60, 1.30, "MAE", "0.4778", "0.4222")
    box(slide, 0.88, 3.15, 5.55, 1.78, "PPL long run on Topical-Chat / USR engagingness\nBaseline prompt -> optimized prompt\nAll reported metrics move in the favorable direction", fill="white", line="blue", size=17, align=PP_ALIGN.LEFT)
    box(slide, 7.00, 3.15, 5.12, 1.78, "Model-confidence feedback can help GEPA sharpen a judge prompt", fill="green_light", line="green", size=21, bold=True)
    textbox(slide, 1.00, 5.65, 11.25, 0.34, "Selected prompt emphasizes specificity, conversational momentum, meaningful follow-up, and concrete detail", size=16, bold=True, color="gray", align=PP_ALIGN.CENTER)
    footer(slide, 9)

    # Slide 10
    slide = prs.slides.add_slide(blank)
    title(slide, SLIDE_TITLES[9])
    add_figure(slide, 0.70, 1.18, 7.55, 3.40, gepa_long_run_metrics, line="blue")
    box(slide, 8.55, 1.20, 3.85, 0.74, "Matched long-run evidence", fill="blue_light", line="blue", size=16, bold=True)
    box(slide, 8.55, 2.12, 3.85, 0.82, "Raw NLA underperformed\nmatched PPL-only", fill="red_light", line="red", size=14.2, bold=True, align=PP_ALIGN.LEFT)
    box(slide, 8.55, 3.14, 3.85, 0.82, "Schema-fixed NLA:\nsmall upward movement", fill="green_light", line="green", size=14.2, bold=True, align=PP_ALIGN.LEFT)
    box(slide, 8.55, 4.16, 3.85, 0.82, "Selected prompt\nstayed unchanged", fill="light_gray", line="mid_gray", size=14.2, bold=True, align=PP_ALIGN.LEFT)
    box(slide, 0.92, 5.00, 3.55, 0.75, "Raw NLA\nPearson 0.5111, MAE 0.6167", fill="red_light", line="red", size=13.4, bold=True)
    box(slide, 4.88, 5.00, 3.55, 0.75, "Schema-fixed\nPearson 0.6812, MAE 0.4944", fill="green_light", line="green", size=13.4, bold=True)
    box(slide, 8.84, 5.00, 3.25, 0.75, "2/60 final-test\npredictions changed", fill="light_gray", line="mid_gray", size=13.4, bold=True)
    textbox(slide, 0.92, 6.05, 11.55, 0.40, "Conclusion: NLA is diagnostic evidence, but not yet a robust prompt-improvement signal", size=17, bold=True, color="blue", align=PP_ALIGN.CENTER)
    footer(slide, 10)

    # Slide 11
    slide = prs.slides.add_slide(blank)
    title(slide, SLIDE_TITLES[10])
    box(slide, 0.82, 1.25, 5.62, 0.55, "What the thesis established", fill="green_light", line="green", size=17, bold=True)
    box(slide, 6.90, 1.25, 5.62, 0.55, "What remains open", fill="blue_light", line="blue", size=17, bold=True)
    achieved = [
        "Controlled corpus for semantic drift",
        "Discrete recovery vs soft-prompt opacity",
        "Matched Latent-GEPA feedback tests",
        "Clearest positive evidence: PPL feedback",
    ]
    future = [
        "NLA token/layer selection, deduplication, compression",
        "Target-score probabilities beyond PPL",
        "Matched SummEval and QAGS matrix",
        "Stronger continuous-prompt readout controls",
    ]
    y_positions = [2.05, 2.85, 3.65, 4.45]
    for y, text in zip(y_positions, achieved):
        box(slide, 1.02, y, 5.15, 0.52, text, fill="white", line="green", size=15.2, bold=True, align=PP_ALIGN.LEFT)
    for y, text in zip(y_positions, future):
        box(slide, 7.10, y, 5.15, 0.52, text, fill="white", line="blue", size=15.2, bold=True, align=PP_ALIGN.LEFT)
    textbox(
        slide,
        1.05,
        5.72,
        11.10,
        0.50,
        "Latent evidence is useful, but its interface to prompt optimization must be designed and validated carefully",
        size=17,
        bold=True,
        color="black",
        align=PP_ALIGN.CENTER,
    )
    footer(slide, 11)

    prs.save(PPTX_PATH)
    return prs


def get_font(size, bold=False):
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def shape_rgb(shape):
    try:
        rgb_value = shape.fill.fore_color.rgb
        if rgb_value:
            return tuple(rgb_value)
    except Exception:
        pass
    return (255, 255, 255)


def line_rgb(shape):
    try:
        rgb_value = shape.line.color.rgb
        if rgb_value:
            return tuple(rgb_value)
    except Exception:
        pass
    return (205, 211, 219)


def shape_text(shape, flatten=False):
    if not hasattr(shape, "text"):
        return ""
    try:
        paragraphs = []
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()
            if text:
                paragraphs.append(text)
        text = "\n".join(paragraphs)
    except Exception:
        text = shape.text
    if flatten:
        return " ".join(text.split())
    return text


def deck_hyperlinks(prs):
    links = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    try:
                        address = run.hyperlink.address
                    except Exception:
                        address = None
                    if address:
                        label = run.text.strip()
                        links[label] = address
    return links


def deck_font_sizes(prs):
    sizes = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size:
                        sizes.add(round(run.font.size.pt, 1))
    return sorted(sizes)


def draw_wrapped(draw, xy, text, width_px, font, fill, max_lines=8):
    if not text:
        return
    avg = max(5, int(font.size * 0.54))
    chars = max(8, int(width_px / avg))
    lines = []
    for paragraph in text.split("\n"):
        lines.extend(wrap(paragraph, width=chars) or [""])
    lines = lines[:max_lines]
    draw.multiline_text(xy, "\n".join(lines), font=font, fill=fill, spacing=3)


def render_preview():
    PREVIEW_DIR.mkdir(exist_ok=True)
    prs = Presentation(PPTX_PATH)
    scale = 1600 / prs.slide_width
    for idx, slide in enumerate(prs.slides, 1):
        img = Image.new("RGB", (1600, 900), "white")
        draw = ImageDraw.Draw(img)
        for shape in slide.shapes:
            try:
                x = int(shape.left * scale)
                y = int(shape.top * scale)
                w = int(shape.width * scale)
                h = int(shape.height * scale)
            except Exception:
                continue
            text = shape_text(shape)
            fill = shape_rgb(shape)
            outline = line_rgb(shape)
            if w < 12 or h < 8:
                continue
            if hasattr(shape, "image"):
                try:
                    source = Image.open(BytesIO(shape.image.blob)).convert("RGB")
                    rendered = source.resize((w, h), Image.Resampling.LANCZOS)
                    img.paste(rendered, (x, y))
                    continue
                except Exception:
                    pass
            if text or shape.shape_type != 17:
                draw.rounded_rectangle([x, y, x + w, y + h], radius=10, fill=fill, outline=outline, width=2)
            if text:
                font_size = 18
                try:
                    first_run = shape.text_frame.paragraphs[0].runs[0]
                    if first_run.font.size:
                        font_size = max(9, min(38, int(first_run.font.size.pt * 1.35)))
                    bold = bool(first_run.font.bold)
                except Exception:
                    bold = False
                font = get_font(font_size, bold=bold)
                draw_wrapped(draw, (x + 12, y + 8), text, max(40, w - 24), font, (18, 24, 32), max_lines=max(1, h // max(14, font_size + 3)))
        img.save(PREVIEW_DIR / f"slide_{idx:02d}.png")


def overlap(a, b):
    ax1, ay1, ax2, ay2 = a
    bx1, by1, bx2, by2 = b
    ix1, iy1 = max(ax1, bx1), max(ay1, by1)
    ix2, iy2 = min(ax2, bx2), min(ay2, by2)
    if ix2 <= ix1 or iy2 <= iy1:
        return 0
    return (ix2 - ix1) * (iy2 - iy1)


def qa_report():
    prs = Presentation(PPTX_PATH)
    issues = []
    titles = []
    required_links = {
        "GEPA: Agrawal et al., 2025": "https://doi.org/10.48550/arXiv.2507.19457",
        "G-EVAL: Liu et al., 2023": "https://doi.org/10.18653/v1/2023.emnlp-main.153",
        "NLA: Fraser-Taliente et al., 2026": "https://transformer-circuits.pub/2026/nla/",
        "SIPIT: Nikolaou et al., 2025": "https://doi.org/10.48550/arXiv.2510.15511",
    }
    links = deck_hyperlinks(prs)
    font_sizes = deck_font_sizes(prs)
    picture_count = sum(1 for slide in prs.slides for shape in slide.shapes if hasattr(shape, "image"))
    for idx, slide in enumerate(prs.slides, 1):
        texts = [shape_text(s, flatten=True) for s in slide.shapes if shape_text(s, flatten=True)]
        if idx > 1 and not any("Simone Mazzacano | Latent-GEPA" in t for t in texts):
            issues.append(f"Slide {idx}: missing footer")
        if idx > 1 and not any(t == str(idx) for t in texts):
            issues.append(f"Slide {idx}: missing slide number")
        if any("bibliography" in t.lower() for t in texts):
            issues.append(f"Slide {idx}: bibliography text detected")
        if any("unibo" in t.lower() for t in texts):
            issues.append(f"Slide {idx}: Unibo text detected")
        if idx > 1:
            titles.append(texts[0] if texts else "")
        boxes = []
        for s in slide.shapes:
            t = shape_text(s, flatten=True)
            if not t or "Simone Mazzacano" in t or t == str(idx) or "Bachelor Thesis Defense" in t:
                continue
            try:
                # Ignore tiny citation links and decorative chips.
                if s.height < Inches(0.25):
                    continue
                boxes.append((s.left, s.top, s.left + s.width, s.top + s.height, t[:40]))
            except Exception:
                continue
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                a = boxes[i]
                b = boxes[j]
                inter = overlap(a[:4], b[:4])
                if inter == 0:
                    continue
                area_a = (a[2] - a[0]) * (a[3] - a[1])
                area_b = (b[2] - b[0]) * (b[3] - b[1])
                if inter / min(area_a, area_b) > 0.08:
                    issues.append(f"Slide {idx}: possible overlap between '{a[4]}' and '{b[4]}'")
    if not 10 <= len(prs.slides) <= 12:
        issues.append(f"Slide count is {len(prs.slides)}, expected 10-12")
    if len(titles) != len(set(titles)):
        issues.append("Duplicate content-slide title detected")
    if picture_count < 2:
        issues.append(f"Only {picture_count} embedded slide images detected, expected at least 2")
    if len(font_sizes) > len(FONT_SCALE):
        issues.append(f"Unexpected font-size count: {font_sizes}")
    for label, address in required_links.items():
        if links.get(label) != address:
            issues.append(f"Missing required paper link: {label}")

    lines = [
        "# Deck QA Report",
        "",
        f"- PPTX: `{PPTX_PATH.name}`",
        f"- Slides: {len(prs.slides)}",
        "- Target range: 10-12 slides",
        "- Format: native `.pptx` generated with `python-pptx`",
        "- Local preview: `preview/slide_XX.png`",
        "",
        "## Checks",
        "",
        f"- Slide count in range: {'yes' if 10 <= len(prs.slides) <= 12 else 'no'}",
        f"- Content titles unique: {'yes' if len(titles) == len(set(titles)) else 'no'}",
        f"- Footer present on content slides: {'yes' if not any('missing footer' in i for i in issues) else 'no'}",
        f"- Slide numbers present on content slides: {'yes' if not any('missing slide number' in i for i in issues) else 'no'}",
        f"- Bibliography slide detected: {'yes' if any('bibliography' in i for i in issues) else 'no'}",
        f"- Unibo logo/text detected: {'yes' if any('Unibo' in i for i in issues) else 'no'}",
        f"- Required paper links present: {'yes' if not any('paper link' in i for i in issues) else 'no'}",
        f"- Embedded thesis/result figures: {picture_count}",
        f"- Font scale used: {', '.join(str(size) for size in font_sizes)}",
        f"- Possible text overlaps detected: {'yes' if any('overlap' in i for i in issues) else 'no'}",
        "",
        "## Issues",
        "",
    ]
    if issues:
        lines.extend([f"- {issue}" for issue in issues])
    else:
        lines.append("- None detected by local structural checks")
    lines.extend(
        [
            "",
            "## Required Manual Check",
            "",
            "- Upload `latent_gepa_discussione.pptx` to PowerPoint for the web",
            "- Confirm the file opens in edit mode",
            "- Check each slide for text wrapping and alignment",
            "- Add or test one comment before sending the link to advisors",
        ]
    )
    QA_PATH.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return issues


def main():
    build_deck()
    render_preview()
    issues = qa_report()
    print(f"Wrote {PPTX_PATH}")
    print(f"Wrote previews in {PREVIEW_DIR}")
    print(f"Wrote {QA_PATH}")
    if issues:
        print("QA issues:")
        for issue in issues:
            print(f"- {issue}")


if __name__ == "__main__":
    main()
